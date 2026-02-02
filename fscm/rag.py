from typing import Optional
from pathlib import Path
import typer
from rich.console import Console
import os
import dotenv
import traceback
import hashlib
from datetime import datetime
dotenv.load_dotenv()


# LangChain Provider-specific Imports
from langchain_chroma import Chroma
from langchain_ollama import OllamaEmbeddings
from langchain_experimental.text_splitter import SemanticChunker
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain_community.document_loaders import (
    DirectoryLoader,
    PyPDFLoader,
    TextLoader,
    Docx2txtLoader,
    JSONLoader,
    CSVLoader,
    BSHTMLLoader,
)
from langchain_core.documents import Document

# Optional imports for additional formats
try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from lxml import etree
    LXML_AVAILABLE = True
except ImportError:
    LXML_AVAILABLE = False
app = typer.Typer(no_args_is_help=True)
console = Console()

#setup chroma db

embedding_function = OllamaEmbeddings(model="embeddinggemma")

'''embedding_function = GoogleGenerativeAIEmbeddings(
    model="models/text-embedding-004",
    google_api_key=os.getenv("GOOGLE_API_KEY")
)'''

# Vector store directory
VECTOR_STORE_DIR = "./vector_store"
Path(VECTOR_STORE_DIR).mkdir(exist_ok=True)
# Initialize project database
db = Chroma(
    embedding_function=embedding_function,
    persist_directory="./rag",
    collection_name="rag_collection"
)

# Initialize global knowledge base
knowledge_db = Chroma(
    embedding_function=embedding_function,
    persist_directory="./rag",
    collection_name="knowledge_base"
)

# Initialize splitter
splitter = SemanticChunker(
    embeddings=embedding_function,
    buffer_size=1,
    breakpoint_threshold_type="percentile", 
    breakpoint_threshold_amount=95
)

# Alternative splitter for structured/mixed content (PDFs, Office docs, data formats)
structured_splitter = RecursiveCharacterTextSplitter(
    chunk_size=1000,
    chunk_overlap=200,
    length_function=len,
    separators=["\n\n", "\n", ". ", " ", ""]
)

# Supported file patterns for import
SUPPORTED_FILE_PATTERNS = [
    '*.md', '*.txt',           # Text
    '*.pdf',                    # PDF
    '*.docx', '*.doc',          # Word
    '*.xlsx', '*.xls',          # Excel
    '*.pptx', '*.ppt',          # PowerPoint
    '*.json', '*.xml',          # Data
    '*.csv',                    # CSV
    '*.html', '*.htm',          # HTML
]

class ExcelLoader:
    """Custom Excel loader using pandas (no unstructured dependency)."""

    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self) -> list:
        if not PANDAS_AVAILABLE:
            raise ImportError("pandas is required for Excel files: pip install pandas openpyxl")

        documents = []
        xlsx = pd.ExcelFile(self.file_path)

        for sheet_name in xlsx.sheet_names:
            df = pd.read_excel(xlsx, sheet_name=sheet_name)
            # Convert DataFrame to readable text
            content = f"Sheet: {sheet_name}\n\n"
            content += df.to_string(index=False)
            documents.append(Document(
                page_content=content,
                metadata={"source": self.file_path, "sheet": sheet_name}
            ))

        return documents


class PowerPointLoader:
    """Custom PowerPoint loader using python-pptx (no unstructured dependency)."""

    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self) -> list:
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx is required for PowerPoint files: pip install python-pptx")

        prs = Presentation(self.file_path)
        documents = []

        for slide_num, slide in enumerate(prs.slides, 1):
            text_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)

            if text_parts:
                content = f"Slide {slide_num}:\n" + "\n".join(text_parts)
                documents.append(Document(
                    page_content=content,
                    metadata={"source": self.file_path, "slide": slide_num}
                ))

        return documents


class XMLLoader:
    """Custom XML loader using lxml (no unstructured dependency)."""

    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self) -> list:
        if not LXML_AVAILABLE:
            # Fallback to standard library
            import xml.etree.ElementTree as ET
            tree = ET.parse(self.file_path)
            root = tree.getroot()

            def extract_text(element, depth=0):
                texts = []
                tag = element.tag.split('}')[-1] if '}' in element.tag else element.tag
                if element.text and element.text.strip():
                    texts.append(f"{'  '*depth}{tag}: {element.text.strip()}")
                for child in element:
                    texts.extend(extract_text(child, depth + 1))
                return texts

            content = "\n".join(extract_text(root))
        else:
            tree = etree.parse(self.file_path)
            root = tree.getroot()

            def extract_text(element, depth=0):
                texts = []
                tag = etree.QName(element).localname if element.tag.startswith('{') else element.tag
                if element.text and element.text.strip():
                    texts.append(f"{'  '*depth}{tag}: {element.text.strip()}")
                for child in element:
                    texts.extend(extract_text(child, depth + 1))
                return texts

            content = "\n".join(extract_text(root))

        return [Document(page_content=content, metadata={"source": self.file_path})]


def get_file_loader(file_path: Path):
    """
    Get the appropriate loader based on file extension.

    Args:
        file_path: Path to the file

    Returns:
        Appropriate LangChain document loader

    Supported formats:
        - Documents: .pdf, .docx, .doc, .txt, .md
        - Spreadsheets: .xlsx, .xls, .csv
        - Presentations: .pptx, .ppt
        - Data: .json, .xml
        - Web: .html, .htm
    """
    file_extension = file_path.suffix.lower()

    # PDF
    if file_extension == '.pdf':
        return PyPDFLoader(str(file_path))

    # Text-based formats
    elif file_extension in ['.md', '.txt']:
        return TextLoader(str(file_path), encoding="utf8")

    # Word documents
    elif file_extension in ['.docx', '.doc']:
        return Docx2txtLoader(str(file_path))

    # Excel spreadsheets
    elif file_extension in ['.xlsx', '.xls']:
        return ExcelLoader(str(file_path))

    # PowerPoint presentations
    elif file_extension in ['.pptx', '.ppt']:
        return PowerPointLoader(str(file_path))

    # JSON files
    elif file_extension == '.json':
        return JSONLoader(
            file_path=str(file_path),
            jq_schema='.',
            text_content=False
        )

    # XML files
    elif file_extension == '.xml':
        return XMLLoader(str(file_path))

    # CSV files
    elif file_extension == '.csv':
        return CSVLoader(str(file_path), encoding="utf8")

    # HTML files
    elif file_extension in ['.html', '.htm']:
        return BSHTMLLoader(str(file_path))

    else:
        raise ValueError(f"Unsupported file type: {file_extension}")

def get_appropriate_splitter(file_path: Path):
    """
    Get the appropriate text splitter based on file type.

    Args:
        file_path: Path to the file

    Returns:
        Appropriate text splitter
    """
    file_extension = file_path.suffix.lower()

    # Use SemanticChunker for pure text formats
    if file_extension in ['.md', '.txt']:
        return splitter

    # Use RecursiveCharacterTextSplitter for structured/mixed content
    # (PDFs, Office docs, data formats)
    else:
        return structured_splitter


def calculate_file_hash(file_path: Path) -> str:
    """
    Calculate MD5 hash of a file.

    Args:
        file_path: Path to the file

    Returns:
        MD5 hash as hexadecimal string
    """
    hasher = hashlib.md5()
    with open(file_path, 'rb') as f:
        for chunk in iter(lambda: f.read(8192), b''):
            hasher.update(chunk)
    return hasher.hexdigest()


def get_existing_file_info(source_file: str, project_id: int) -> Optional[dict]:
    """
    Check if file already exists in database and return its info.

    Args:
        source_file: Name of the source file
        project_id: Project ID to check

    Returns:
        dict with 'exists', 'file_hash', 'chunk_ids' if found, None otherwise
    """
    try:
        results = db.get(
            where={"$and": [
                {"source_file": source_file},
                {"project_id": project_id}
            ]},
            include=["metadatas"]
        )
        if results['ids']:
            # Get file_hash from first chunk's metadata
            file_hash = results['metadatas'][0].get('file_hash') if results['metadatas'] else None
            return {
                'exists': True,
                'file_hash': file_hash,
                'chunk_ids': results['ids'],
                'chunk_count': len(results['ids'])
            }
        return None
    except Exception as e:
        print(f"  ⚠️  Error checking existing file: {str(e)}")
        return None


def delete_file_chunks(source_file: str, project_id: int) -> int:
    """
    Delete all chunks for a specific file and project.

    Args:
        source_file: Name of the source file
        project_id: Project ID

    Returns:
        Number of chunks deleted
    """
    try:
        where_filter = {"$and": [
            {"source_file": source_file},
            {"project_id": project_id}
        ]}

        # Get count before deletion
        before_result = db.get(where=where_filter)
        count = len(before_result['ids'])

        if count > 0:
            collection = db._collection
            collection.delete(where=where_filter)

        return count
    except Exception as e:
        print(f"  ⚠️  Error deleting chunks: {str(e)}")
        return 0


def process_single_file(file_path: Path, project_id: int, force: bool = False) -> dict:
    """
    Process a single file and add it to the vector database.
    Supports incremental updates: skips unchanged files, updates modified files.

    Supported formats: PDF, Word, Excel, PowerPoint, TXT, MD, JSON, XML, CSV, HTML

    Args:
        file_path: Path to the file to process
        project_id: Required project ID for the document
        force: Force re-import even if file is unchanged

    Returns:
        dict: Processing result with 'status' (NEW, UPDATED, SKIPPED, FAILED),
              'chunks_added', 'chunks_removed'
    """
    result = {
        'status': 'FAILED',
        'chunks_added': 0,
        'chunks_removed': 0
    }

    try:
        # 1. Calculate file hash
        new_hash = calculate_file_hash(file_path)
        file_size = file_path.stat().st_size

        # 2. Check if file exists in database
        existing = get_existing_file_info(file_path.name, project_id)

        if existing:
            if existing['file_hash'] == new_hash and not force:
                # File unchanged - skip
                result['status'] = 'SKIPPED'
                return result
            else:
                # File changed or force mode - delete old chunks first
                deleted_count = delete_file_chunks(file_path.name, project_id)
                result['chunks_removed'] = deleted_count
                result['status'] = 'UPDATED'
        else:
            result['status'] = 'NEW'

        # 3. Get appropriate loader for file type
        try:
            loader = get_file_loader(file_path)
        except ValueError as e:
            print(f"  ⚠️  {str(e)}")
            result['status'] = 'FAILED'
            return result

        # Load document
        docs = loader.load()

        if not docs:
            print(f"  ⚠️  No content found in {file_path.name}")
            result['status'] = 'FAILED'
            return result

        # For PDFs, combine all pages into content info
        total_content = ""
        if file_path.suffix.lower() == '.pdf':
            total_content = "\n".join([doc.page_content for doc in docs])
        else:
            total_content = docs[0].page_content

        # Get appropriate splitter for file type
        current_splitter = get_appropriate_splitter(file_path)

        # Split document into chunks
        chunks = current_splitter.split_documents(docs)

        if not chunks:
            print(f"  ⚠️  No chunks created from {file_path.name}")
            result['status'] = 'FAILED'
            return result

        # Add enhanced metadata for tracking (including new hash fields)
        for i, chunk in enumerate(chunks):
            chunk.metadata.update({
                'source_file': file_path.name,
                'file_path': str(file_path),
                'file_type': file_path.suffix.lower(),
                'processed_at': datetime.now().isoformat(),
                'chunk_count': len(chunks),
                'chunk_index': i,
                'splitter_used': current_splitter.__class__.__name__,
                'project_id': project_id,
                'file_hash': new_hash,
                'file_size': file_size
            })

            # For PDFs, add page information if available
            if file_path.suffix.lower() == '.pdf' and 'page' in chunk.metadata:
                chunk.metadata['source_page'] = chunk.metadata.get('page', 'unknown')

        # Add chunks to database in batches (Chroma max batch size is 5461)
        BATCH_SIZE = 5000
        total_chunks = len(chunks)

        if total_chunks > BATCH_SIZE:
            print(f"  📦 Large file: splitting {total_chunks} chunks into batches of {BATCH_SIZE}")

        for i in range(0, total_chunks, BATCH_SIZE):
            batch = chunks[i:i + BATCH_SIZE]
            db.add_documents(batch)
            if total_chunks > BATCH_SIZE:
                batch_num = (i // BATCH_SIZE) + 1
                total_batches = (total_chunks + BATCH_SIZE - 1) // BATCH_SIZE
                print(f"  📦 Batch {batch_num}/{total_batches} ({len(batch)} chunks)")

        result['chunks_added'] = len(chunks)

        return result

    except Exception as e:
        print(f"  ❌ Error processing {file_path.name}: {str(e)}")
        print(f"  📝 Traceback: {traceback.format_exc()}")
        result['status'] = 'FAILED'
        return result

def scan_and_process_directory(directory_path: str, project_id: int, file_types: list = None, force: bool = False) -> dict:
    """
    Scan directory and process all supported files individually.
    Supports incremental updates: skips unchanged files, updates modified files.

    Args:
        directory_path: Path to directory to scan
        project_id: Required project ID for all processed documents
        file_types: List of file extensions to process (default: SUPPORTED_FILE_PATTERNS)
        force: Force re-import all files even if unchanged

    Returns:
        dict: Processing statistics
    """
    if file_types is None:
        file_types = SUPPORTED_FILE_PATTERNS

    stats = {
        'total_files': 0,
        'new_files': 0,
        'updated_files': 0,
        'skipped_files': 0,
        'failed_files': 0,
        'new_chunks': 0,
        'updated_chunks': 0,
        'removed_chunks': 0,
        'failed_file_list': [],
        'file_type_stats': {},
        'file_results': []
    }

    # Find all supported files
    directory = Path(directory_path)
    if not directory.exists():
        print(f"❌ Directory does not exist: {directory_path}")
        return stats

    # Get all supported files recursively
    all_files = []
    for file_pattern in file_types:
        found_files = list(directory.rglob(file_pattern))
        all_files.extend(found_files)

        # Track file type statistics
        file_ext = file_pattern.replace('*', '').upper()
        stats['file_type_stats'][file_ext] = len(found_files)

    # Remove duplicates (in case of overlapping patterns)
    all_files = list(set(all_files))
    # Sort for consistent output
    all_files.sort(key=lambda x: x.name)
    stats['total_files'] = len(all_files)

    print(f"\n🔍 Found {len(all_files)} files to process:")
    for file_ext, count in stats['file_type_stats'].items():
        print(f"  - {file_ext} files: {count}")

    if not all_files:
        print("⚠️  No supported files found in directory")
        return stats

    print(f"\nProcessing: {directory_path}/")

    # Process each file individually
    for file_path in all_files:
        result = process_single_file(file_path, project_id=project_id, force=force)

        # Build status indicator
        status = result['status']
        chunks_info = ""

        if status == 'NEW':
            stats['new_files'] += 1
            stats['new_chunks'] += result['chunks_added']
            chunks_info = f"+{result['chunks_added']} chunks"
            status_display = "[NEW]"
        elif status == 'UPDATED':
            stats['updated_files'] += 1
            stats['updated_chunks'] += result['chunks_added']
            stats['removed_chunks'] += result['chunks_removed']
            chunks_info = f"-{result['chunks_removed']}/+{result['chunks_added']} chunks"
            status_display = "[UPDATED]"
        elif status == 'SKIPPED':
            stats['skipped_files'] += 1
            chunks_info = "unchanged"
            status_display = "[SKIPPED]"
        else:  # FAILED
            stats['failed_files'] += 1
            stats['failed_file_list'].append(str(file_path))
            chunks_info = "error"
            status_display = "[FAILED]"

        # Store result for potential later use
        stats['file_results'].append({
            'file': file_path.name,
            'status': status,
            'chunks_added': result['chunks_added'],
            'chunks_removed': result['chunks_removed']
        })

        # Print tree-like output
        prefix = "├──" if file_path != all_files[-1] else "└──"
        print(f"{prefix} {file_path.name:<30} {status_display:<12} {chunks_info}")

    return stats


def preview_deletion(where_filter: dict) -> dict:
    """
    Preview what documents will be deleted without actually deleting them.

    Args:
        where_filter: Chroma where filter dict

    Returns:
        dict: Preview statistics including count, affected files, project IDs
    """
    try:
        # Query database with filter
        result = db.get(where=where_filter)

        count = len(result['ids'])

        if count == 0:
            return {
                'count': 0,
                'source_files': set(),
                'project_ids': set(),
                'sample_metadata': None
            }

        # Extract metadata
        metadatas = result['metadatas']
        source_files = set()
        project_ids = set()

        for metadata in metadatas:
            if 'source_file' in metadata:
                source_files.add(metadata['source_file'])
            if 'project_id' in metadata:
                project_ids.add(metadata['project_id'])

        # Get first chunk as sample
        sample_metadata = metadatas[0] if metadatas else None

        return {
            'count': count,
            'source_files': source_files,
            'project_ids': project_ids,
            'sample_metadata': sample_metadata
        }
    except Exception as e:
        print(f"❌ Error previewing deletion: {str(e)}")
        return {
            'count': 0,
            'source_files': set(),
            'project_ids': set(),
            'sample_metadata': None
        }


def execute_deletion(where_filter: dict) -> int:
    """
    Execute the deletion and return count of deleted documents.

    Args:
        where_filter: Chroma where filter dict

    Returns:
        int: Number of documents deleted
    """
    try:
        # Get before count
        before_count = len(db.get()['ids'])

        # Execute deletion using native Chroma API
        collection = db._collection
        collection.delete(where=where_filter)

        # Get after count
        after_count = len(db.get()['ids'])

        return before_count - after_count
    except Exception as e:
        print(f"❌ Error executing deletion: {str(e)}")
        print(f"  📝 Traceback: {traceback.format_exc()}")
        return 0


def confirm_deletion(count: int, mode: str) -> bool:
    """
    Display warning and get user confirmation for deletion.

    Args:
        count: Number of documents to be deleted
        mode: Deletion mode description (e.g., "file: report.pdf")

    Returns:
        bool: True if user confirms, False otherwise
    """
    console.print(f"\n[bold red]⚠️  WARNING[/bold red]")
    console.print(f"This will permanently delete [bold]{count}[/bold] documents ({mode}).")
    console.print("This action cannot be undone.\n")

    return typer.confirm("Do you want to proceed?", default=False)


@app.command()
def addfiles(
    folder: str = typer.Option(..., "--folder", "-f", help="Folder to Import"),
    project_id: int = typer.Option(..., "--projekt_id", "-p", help="Projekt ID (required)"),
    force: bool = typer.Option(False, "--force", "-F", help="Force re-import all files (even unchanged)"),
):
    """
    Import files from a folder into the RAG vector database.
    Supports incremental updates: skips unchanged files, updates modified files.

    Supported formats:
    - Documents: PDF, Word (.docx/.doc), Text (.txt), Markdown (.md)
    - Spreadsheets: Excel (.xlsx/.xls), CSV
    - Presentations: PowerPoint (.pptx/.ppt)
    - Data: JSON, XML
    - Web: HTML
    """
    print("🚀 Starting file import with change detection...")
    print("📋 Supported formats: PDF, Word, Excel, PowerPoint, TXT, MD, JSON, XML, CSV, HTML")

    if force:
        print("⚠️  Force mode: All files will be re-imported")

    # Check if database already has content
    existing_docs = db.get()
    print(f"📊 Database currently contains {len(existing_docs['ids'])} documents")

    # Process all supported files
    processing_stats = scan_and_process_directory(
        folder,
        project_id,
        file_types=SUPPORTED_FILE_PATTERNS,
        force=force
    )

    # Print summary
    print(f"\n" + "="*50)
    print("📈 SUMMARY")
    print("="*50)

    # New summary format
    new_chunks = processing_stats['new_chunks']
    updated_chunks = processing_stats['updated_chunks']

    print(f"  New files:     {processing_stats['new_files']} ({new_chunks} chunks)")
    print(f"  Updated files: {processing_stats['updated_files']} ({updated_chunks} chunks)")
    print(f"  Skipped:       {processing_stats['skipped_files']}")
    print(f"  Failed:        {processing_stats['failed_files']}")

    if processing_stats['removed_chunks'] > 0:
        print(f"\n  Chunks removed (from updates): {processing_stats['removed_chunks']}")

    # Final database size
    final_count = len(db.get()['ids'])
    print(f"\n📊 Final database size: {final_count} documents")

    if processing_stats['failed_file_list']:
        print(f"\n❌ Failed files:")
        for failed_file in processing_stats['failed_file_list']:
            print(f"  - {failed_file}")


@app.command()
def delete(
    source_file: str = typer.Option(None, "--file", "-f", help="Delete all chunks from specific source file"),
    project_id: int = typer.Option(None, "--project", "-p", help="Delete all documents from specific project"),
    dry_run: bool = typer.Option(False, "--dry-run", "-d", help="Preview deletion without executing"),
    yes: bool = typer.Option(False, "--yes", "-y", help="Skip confirmation prompt"),
):
    """Delete documents from the RAG vector database."""

    # Phase 1: Input Validation
    if source_file is None and project_id is None:
        console.print("[bold red]❌ Error:[/bold red] You must specify either --file or --project")
        console.print("Usage: delete --file <filename>  OR  delete --project <id>")
        raise typer.Exit(code=1)

    if source_file is not None and project_id is not None:
        console.print("[bold red]❌ Error:[/bold red] You cannot specify both --file and --project")
        console.print("Usage: delete --file <filename>  OR  delete --project <id>")
        raise typer.Exit(code=1)

    # Phase 2: Build where filter
    if source_file is not None:
        where_filter = {"source_file": source_file}
        mode_description = f"file: {source_file}"
        mode_type = "source file"
    else:
        where_filter = {"project_id": project_id}
        mode_description = f"project: {project_id}"
        mode_type = "project"

    # Show current database status
    try:
        total_docs = len(db.get()['ids'])
        console.print(f"\n📊 Database currently contains [bold]{total_docs}[/bold] documents")
    except Exception as e:
        console.print(f"[bold red]❌ Error accessing database:[/bold red] {str(e)}")
        raise typer.Exit(code=1)

    # Phase 3: Preview deletion
    console.print(f"\n🔍 Searching for documents to delete...")
    preview_data = preview_deletion(where_filter)

    if preview_data['count'] == 0:
        console.print(f"[yellow]ℹ️  No documents found matching {mode_description}[/yellow]")
        console.print("Nothing to delete.")
        raise typer.Exit(code=0)

    # Display preview
    console.print(f"\n[bold]🗑️  DELETION PREVIEW[/bold]")
    console.print("=" * 50)
    console.print(f"Mode: Delete by {mode_type}")
    if source_file:
        console.print(f"File: [cyan]{source_file}[/cyan]")
    else:
        console.print(f"Project ID: [cyan]{project_id}[/cyan]")

    console.print(f"Documents to delete: [bold red]{preview_data['count']}[/bold red]")

    # Show affected files or projects
    if project_id is not None and preview_data['source_files']:
        console.print(f"\nAffected files ({len(preview_data['source_files'])}):")
        for file in sorted(preview_data['source_files']):
            console.print(f"  - {file}")

    if source_file is not None and preview_data['project_ids']:
        project_list = [str(p) if p is not None else "None" for p in preview_data['project_ids']]
        console.print(f"\nAffected projects: {', '.join(sorted(project_list))}")

    # Show sample metadata
    if preview_data['sample_metadata']:
        sample = preview_data['sample_metadata']
        console.print(f"\nSample metadata:")
        console.print(f"  - File: {sample.get('source_file', 'Unknown')}")
        console.print(f"  - Project: {sample.get('project_id', 'N/A')}")
        console.print(f"  - Type: {sample.get('file_type', 'Unknown')}")
        console.print(f"  - Page: {sample.get('source_page', 'N/A')}")

    # Phase 4: Handle dry-run
    if dry_run:
        console.print(f"\n[bold yellow]🔍 DRY RUN - No changes made[/bold yellow]")
        raise typer.Exit(code=0)

    # Phase 5: Confirmation
    if not yes:
        if not confirm_deletion(preview_data['count'], mode_description):
            console.print("\n[yellow]❌ Deletion cancelled[/yellow]")
            raise typer.Exit(code=0)

    # Phase 6: Execute deletion
    console.print(f"\n🗑️  Deleting documents...")
    deleted_count = execute_deletion(where_filter)

    # Phase 7: Display summary
    if deleted_count == 0:
        console.print("[yellow]⚠️  No documents were deleted[/yellow]")
    elif deleted_count != preview_data['count']:
        console.print(f"[yellow]⚠️  Warning: Expected to delete {preview_data['count']} documents, but deleted {deleted_count}[/yellow]")
    else:
        console.print(f"[bold green]✅ Successfully deleted {deleted_count} documents[/bold green]")

    # Show final database size
    try:
        final_count = len(db.get()['ids'])
        console.print(f"📊 Database now contains [bold]{final_count}[/bold] documents")
    except Exception as e:
        console.print(f"[yellow]⚠️  Could not verify final database size: {str(e)}[/yellow]")

    console.print("\n🎉 Deletion complete!")


# ============================================================================
# GLOBAL KNOWLEDGE BASE COMMANDS
# ============================================================================

def process_single_file_to_knowledge_db(file_path: Path, force: bool = False) -> dict:
    """
    Process a single file and add it to the GLOBAL knowledge base.
    Similar to process_single_file but uses knowledge_db instead of db.

    Args:
        file_path: Path to the file to process
        force: Force re-import even if file is unchanged

    Returns:
        dict: Processing result with 'status', 'chunks_added', 'chunks_removed'
    """
    result = {
        'status': 'FAILED',
        'chunks_added': 0,
        'chunks_removed': 0
    }

    try:
        # 1. Calculate file hash
        new_hash = calculate_file_hash(file_path)
        file_size = file_path.stat().st_size

        # 2. Check if file exists in knowledge database
        try:
            existing = knowledge_db.get(
                where={"source_file": file_path.name},
                include=["metadatas"]
            )
            if existing['ids']:
                file_hash = existing['metadatas'][0].get('file_hash') if existing['metadatas'] else None
                if file_hash == new_hash and not force:
                    result['status'] = 'SKIPPED'
                    return result
                else:
                    # Delete old chunks
                    collection = knowledge_db._collection
                    collection.delete(where={"source_file": file_path.name})
                    result['chunks_removed'] = len(existing['ids'])
                    result['status'] = 'UPDATED'
            else:
                result['status'] = 'NEW'
        except Exception:
            result['status'] = 'NEW'

        # 3. Get appropriate loader
        try:
            loader = get_file_loader(file_path)
        except ValueError as e:
            print(f"  ⚠️  {str(e)}")
            result['status'] = 'FAILED'
            return result

        # Load document
        docs = loader.load()
        if not docs:
            print(f"  ⚠️  No content found in {file_path.name}")
            result['status'] = 'FAILED'
            return result

        # Get appropriate splitter
        current_splitter = get_appropriate_splitter(file_path)

        # Split document into chunks
        chunks = current_splitter.split_documents(docs)
        if not chunks:
            print(f"  ⚠️  No chunks created from {file_path.name}")
            result['status'] = 'FAILED'
            return result

        # Add metadata (NO project_id for global knowledge base)
        for i, chunk in enumerate(chunks):
            chunk.metadata.update({
                'source_file': file_path.name,
                'file_path': str(file_path),
                'file_type': file_path.suffix.lower(),
                'processed_at': datetime.now().isoformat(),
                'chunk_count': len(chunks),
                'chunk_index': i,
                'splitter_used': current_splitter.__class__.__name__,
                'file_hash': new_hash,
                'file_size': file_size,
                'knowledge_base': True  # Mark as global knowledge
            })

            if file_path.suffix.lower() == '.pdf' and 'page' in chunk.metadata:
                chunk.metadata['source_page'] = chunk.metadata.get('page', 'unknown')

        # Add to KNOWLEDGE database in batches (Chroma max batch size is 5461)
        BATCH_SIZE = 5000
        total_chunks = len(chunks)

        if total_chunks > BATCH_SIZE:
            print(f"  📦 Large file: splitting {total_chunks} chunks into batches of {BATCH_SIZE}")

        for i in range(0, total_chunks, BATCH_SIZE):
            batch = chunks[i:i + BATCH_SIZE]
            knowledge_db.add_documents(batch)
            if total_chunks > BATCH_SIZE:
                batch_num = (i // BATCH_SIZE) + 1
                total_batches = (total_chunks + BATCH_SIZE - 1) // BATCH_SIZE
                print(f"  📦 Batch {batch_num}/{total_batches} ({len(batch)} chunks)")

        result['chunks_added'] = len(chunks)

        return result

    except Exception as e:
        print(f"  ❌ Error processing {file_path.name}: {str(e)}")
        result['status'] = 'FAILED'
        return result


@app.command("add-knowledge")
def add_knowledge(
    folder: str = typer.Option(..., "--folder", "-f", help="Folder to import into knowledge base"),
    force: bool = typer.Option(False, "--force", "-F", help="Force re-import all files"),
):
    """
    Import files into the GLOBAL knowledge base (not project-specific).
    Use this for D365 documentation, best practices, general references.

    The global knowledge base is automatically searched alongside project data.
    """
    print("🚀 Importing into GLOBAL knowledge base...")
    print("📋 Supported formats: PDF, Word, Excel, PowerPoint, TXT, MD, JSON, XML, CSV, HTML")

    if force:
        print("⚠️  Force mode: All files will be re-imported")

    # Check current knowledge base size
    try:
        existing_docs = knowledge_db.get()
        print(f"📊 Knowledge base currently contains {len(existing_docs['ids'])} documents")
    except Exception:
        print("📊 Knowledge base is empty")

    # Find files
    directory = Path(folder)
    if not directory.exists():
        console.print(f"[bold red]Error:[/bold red] Directory not found: {folder}")
        raise typer.Exit(code=1)

    all_files = []
    for pattern in SUPPORTED_FILE_PATTERNS:
        all_files.extend(list(directory.rglob(pattern)))
    all_files = sorted(set(all_files), key=lambda x: x.name)

    if not all_files:
        print("⚠️  No supported files found")
        return

    print(f"\n🔍 Found {len(all_files)} files")
    print(f"\nProcessing: {folder}/")

    # Process stats
    stats = {'new': 0, 'updated': 0, 'skipped': 0, 'failed': 0, 'chunks': 0}

    for file_path in all_files:
        result = process_single_file_to_knowledge_db(file_path, force=force)

        status = result['status']
        if status == 'NEW':
            stats['new'] += 1
            stats['chunks'] += result['chunks_added']
            indicator = "[NEW]"
        elif status == 'UPDATED':
            stats['updated'] += 1
            stats['chunks'] += result['chunks_added']
            indicator = "[UPDATED]"
        elif status == 'SKIPPED':
            stats['skipped'] += 1
            indicator = "[SKIPPED]"
        else:
            stats['failed'] += 1
            indicator = "[FAILED]"

        prefix = "├──" if file_path != all_files[-1] else "└──"
        print(f"{prefix} {file_path.name:<30} {indicator}")

    # Summary
    print(f"\n" + "="*50)
    print("📈 SUMMARY")
    print("="*50)
    print(f"  New files:     {stats['new']}")
    print(f"  Updated files: {stats['updated']}")
    print(f"  Skipped:       {stats['skipped']}")
    print(f"  Failed:        {stats['failed']}")
    print(f"  Total chunks:  {stats['chunks']}")

    try:
        final_count = len(knowledge_db.get()['ids'])
        print(f"\n📊 Knowledge base now contains {final_count} documents")
    except Exception:
        pass


@app.command("list-knowledge")
def list_knowledge():
    """List all files in the global knowledge base."""
    try:
        all_docs = knowledge_db.get(include=["metadatas"])

        if not all_docs['ids']:
            console.print("[yellow]Knowledge base is empty[/yellow]")
            return

        # Group by source file
        files = {}
        for metadata in all_docs['metadatas']:
            source = metadata.get('source_file', 'Unknown')
            if source not in files:
                files[source] = {
                    'chunks': 0,
                    'file_type': metadata.get('file_type', '?'),
                    'processed_at': metadata.get('processed_at', '?')
                }
            files[source]['chunks'] += 1

        console.print(f"\n[bold]Global Knowledge Base[/bold] ({len(all_docs['ids'])} total chunks)\n")

        for source, info in sorted(files.items()):
            console.print(f"  {source:<40} {info['file_type']:<8} {info['chunks']} chunks")

        console.print(f"\n[dim]Total: {len(files)} files[/dim]")

    except Exception as e:
        console.print(f"[bold red]Error:[/bold red] {str(e)}")


@app.command("delete-knowledge")
def delete_knowledge(
    source_file: str = typer.Option(None, "--file", "-f", help="Delete specific file from knowledge base"),
    all_files: bool = typer.Option(False, "--all", "-a", help="Delete ALL files from knowledge base"),
    yes: bool = typer.Option(False, "--yes", "-y", help="Skip confirmation"),
):
    """Delete files from the global knowledge base."""
    if not source_file and not all_files:
        console.print("[bold red]Error:[/bold red] Specify --file or --all")
        raise typer.Exit(code=1)

    try:
        if all_files:
            # Delete everything
            all_docs = knowledge_db.get()
            count = len(all_docs['ids'])

            if count == 0:
                console.print("[yellow]Knowledge base is already empty[/yellow]")
                return

            if not yes:
                console.print(f"[bold red]⚠️  This will delete ALL {count} documents from the knowledge base[/bold red]")
                if not typer.confirm("Continue?", default=False):
                    console.print("[yellow]Cancelled[/yellow]")
                    return

            # Delete all
            collection = knowledge_db._collection
            collection.delete(where={})
            console.print(f"[bold green]✅ Deleted {count} documents[/bold green]")

        else:
            # Delete specific file
            existing = knowledge_db.get(where={"source_file": source_file})
            count = len(existing['ids'])

            if count == 0:
                console.print(f"[yellow]File '{source_file}' not found in knowledge base[/yellow]")
                return

            if not yes:
                console.print(f"[bold]Delete '{source_file}' ({count} chunks)?[/bold]")
                if not typer.confirm("Continue?", default=False):
                    console.print("[yellow]Cancelled[/yellow]")
                    return

            collection = knowledge_db._collection
            collection.delete(where={"source_file": source_file})
            console.print(f"[bold green]✅ Deleted {count} chunks from '{source_file}'[/bold green]")

        # Show final size
        final_count = len(knowledge_db.get()['ids'])
        console.print(f"📊 Knowledge base now contains {final_count} documents")

    except Exception as e:
        console.print(f"[bold red]Error:[/bold red] {str(e)}")
        raise typer.Exit(code=1)

